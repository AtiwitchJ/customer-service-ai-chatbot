"""
Step A: Universal Extractor (PDF, PPTX, DOCX, XLSX, Image)
==========================================================
Unified extractor logic for the RAG pipeline / 
ตรรกะการดึงข้อมูลแบบรวมศูนย์สำหรับ RAG pipeline
"""

import io
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument
from PIL import Image as PilImage
from dataclasses import dataclass, field
from typing import List, Optional, Tuple
import os
import time
from datetime import datetime
import zipfile
from odf import text, teletype
from odf.opendocument import load as load_odf
from fontTools.ttLib import TTFont


@dataclass
class DocumentElement:
    """Element extracted from document (Text or Image) / Element ที่ดึงมาจากเอกสาร"""
    page_number: int                                   # Page number / เลขหน้า
    position_in_page: int                              # Order within page / ลำดับในหน้า
    y_position: float                                  # Approx Y position / ตำแหน่ง Y โดยประมาณ
    element_type: str                                  # 'text' or 'image' / ประเภท
    content: str = ""                                  # Text content / เนื้อหาข้อความ
    image_bytes: Optional[bytes] = None                # Image bytes / ข้อมูลรูปภาพ
    image_ext: str = ""                                # Extension: png, jpg / นามสกุล
    bbox: Tuple[float, float, float, float] = field(default_factory=lambda: (0, 0, 0, 0))  # Bounding box
    
    def __repr__(self):
        if self.element_type == 'text':
            preview = self.content[:50] + "..." if len(self.content) > 50 else self.content
            return f"TextElement(page={self.page_number}, pos={self.position_in_page}, text='{preview}')"
        return f"ImageElement(page={self.page_number}, pos={self.position_in_page}, ext={self.image_ext})"


@dataclass
class PageContent:
    """Page-level content container / ตัวเก็บเนื้อหาระดับหน้า"""
    page_number: int                                   # Page number / เลขหน้า
    text_blocks: List[DocumentElement] = field(default_factory=list)   # Text blocks / บล็อกข้อความ
    images: List[DocumentElement] = field(default_factory=list)        # Images / รูปภาพ
    full_page_image: Optional[bytes] = None            # For visual OCR / สำหรับ visual OCR

    def get_native_text(self) -> str:
        """Get native text from all text blocks / ดึงข้อความ native จากทุกบล็อก"""
        texts = [elem.content for elem in self.text_blocks if elem.content]
        return "\n".join(texts)
    
    def get_all_elements_sorted(self) -> List[DocumentElement]:
        """Get all elements sorted by position / ดึงทุก element เรียงตามตำแหน่ง"""
        all_elems = self.text_blocks + self.images
        return sorted(all_elems, key=lambda x: (x.y_position, x.position_in_page))


def extract_pages_from_bytes(file_bytes: bytes, filename: str, render_pages: bool = True) -> List[PageContent]:
    """
    Main Entry Point: Dispatch to specific extractor based on extension / 
    จุดเริ่มต้นหลัก: ส่งต่อไปยัง extractor เฉพาะตามนามสกุล
    """
    ext = os.path.splitext(filename)[1].lower()
    print(f"[{datetime.now().isoformat()}] [Extractor] Start extracting {filename} (Ext: {ext}). Size: {len(file_bytes)} bytes")
    
    if ext == '.pdf':
        return _extract_pdf(file_bytes, filename, render_pages)
    elif ext in ['.pptx', '.ppt']:
        return _extract_pptx(file_bytes, filename)
    elif ext in ['.docx', '.doc']:
        return _extract_docx(file_bytes, filename)
    elif ext in ['.xlsx', '.xls', '.csv']:
        return _extract_excel_csv(file_bytes, filename, ext)
    elif ext in ['.jpg', '.jpeg', '.png', '.webp']:
        return _extract_image_file(file_bytes, filename)
    elif ext in ['.txt', '.json', '.xml', '.html', '.rtf', '.md', '.log']:
        return _extract_text_file(file_bytes, filename)
    elif ext in ['.odp', '.ods', '.odt']:
        return _extract_odf(file_bytes, filename)
    elif ext == '.zip':
        return _extract_zip(file_bytes, filename)
    elif ext == '.ttf':
        return _extract_ttf(file_bytes, filename)
    else:
        # Fallback for others / กรณีไม่รองรับ
        print(f"⚠️ Format {ext} allowed but no specific extractor. Skipping text extraction.")
        return []


# ============================================================================
# 1. PDF Extractor (Original Logic)
# 1. PDF Extractor (ตรรกะเดิม)
# ============================================================================
def _extract_pdf(pdf_bytes: bytes, filename: str, render_pages: bool = True, dpi: int = 150) -> List[PageContent]:
    start_pdf = time.time()
    pages: List[PageContent] = []
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    for page_num, page in enumerate(doc, start=1):
        page_content = PageContent(page_number=page_num)
        page_elements: List[Tuple[float, DocumentElement]] = []
        
        # --- Extract Text / ดึงข้อความ ---
        text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
        for block in text_dict.get("blocks", []):
            if block.get("type") == 0:
                y_pos = block.get("bbox", [0, 0, 0, 0])[1]
                lines_text = []
                for line in block.get("lines", []):
                    spans_text = "".join(span.get("text", "") for span in line.get("spans", []))
                    lines_text.append(spans_text)
                block_text = "\n".join(lines_text).strip()
                if block_text:
                    elem = DocumentElement(
                        page_number=page_num, position_in_page=0, y_position=y_pos,
                        element_type='text', content=block_text, bbox=tuple(block.get("bbox", [0, 0, 0, 0]))
                    )
                    page_elements.append((y_pos, elem))
        
        # --- Extract Images / ดึงรูปภาพ ---
        image_list = page.get_images(full=True)
        for img_index, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                img_rects = page.get_image_rects(xref)
                if img_rects:
                    rect = img_rects[0]
                    y_pos = rect.y0
                    bbox = (rect.x0, rect.y0, rect.x1, rect.y1)
                else:
                    y_pos = img_index * 100
                    bbox = (0, 0, 0, 0)
                
                elem = DocumentElement(
                    page_number=page_num, position_in_page=0, y_position=y_pos,
                    element_type='image', image_bytes=image_bytes, image_ext=image_ext, bbox=bbox
                )
                page_elements.append((y_pos, elem))
            except Exception as e:
                print(f"⚠️ Error extracting image {xref}: {e}")
                continue
        
        # Sort elements by Y position / เรียง element ตามตำแหน่ง Y
        page_elements.sort(key=lambda x: x[0])
        for idx, (_, elem) in enumerate(page_elements):
            elem.position_in_page = idx
            if elem.element_type == 'text':
                page_content.text_blocks.append(elem)
            else:
                page_content.images.append(elem)
        
        # --- Full Page Image / รูปภาพทั้งหน้า ---
        if render_pages:
            try:
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                
                # Convert PyMuPDF pixmap to PNG bytes / แปลง pixmap เป็น PNG bytes
                img_data = pix.tobytes("png")
                pil_image = PilImage.open(io.BytesIO(img_data))
                
                # Convert back to bytes / แปลงกลับเป็น bytes
                output_buffer = io.BytesIO()
                pil_image.save(output_buffer, format='PNG', optimize=True)
                page_content.full_page_image = output_buffer.getvalue()
                
            except Exception as e:
                print(f"⚠️ Error rendering page: {e}")
        
        pages.append(page_content)
    
    doc.close()
    duration = time.time() - start_pdf
    print(f"[{datetime.now().isoformat()}] [Extractor:PDF] Processed {len(pages)} pages in {duration:.2f}s")
    return pages


# ============================================================================
# 2. PPTX Extractor
# 2. PPTX Extractor
# ============================================================================
def _extract_pptx(file_bytes: bytes, filename: str) -> List[PageContent]:
    pages: List[PageContent] = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide_idx, slide in enumerate(prs.slides, start=1):
            page_content = PageContent(page_number=slide_idx)
            
            # Sort shapes by position to mimic reading order / เรียง shapes ตามตำแหน่ง
            shapes = sorted(slide.shapes, key=lambda s: (s.top if hasattr(s, 'top') else 0, s.left if hasattr(s, 'left') else 0))
            
            pos_counter = 0
            for shape in shapes:
                # Text / ข้อความ
                if hasattr(shape, "text") and shape.text.strip():
                    page_content.text_blocks.append(DocumentElement(
                        page_number=slide_idx,
                        position_in_page=pos_counter,
                        y_position=float(shape.top) if hasattr(shape, 'top') else 0.0,
                        element_type='text',
                        content=shape.text.strip()
                    ))
                    pos_counter += 1
                
                # Image / รูปภาพ (shape_type 13 is PICTURE)
                if shape.shape_type == 13: 
                    try:
                        image_blob = shape.image.blob
                        image_ext = shape.image.ext
                        page_content.images.append(DocumentElement(
                            page_number=slide_idx,
                            position_in_page=pos_counter,
                            y_position=float(shape.top) if hasattr(shape, 'top') else 0.0,
                            element_type='image',
                            image_bytes=image_blob,
                            image_ext=image_ext
                        ))
                        pos_counter += 1
                    except Exception as e:
                        print(f"⚠️ Error extracting PPTX image: {e}")

            pages.append(page_content)
        print(f"✅ Extracted {len(pages)} slides from PPTX")
    except Exception as e:
        print(f"❌ Error extracting PPTX: {e}")
    return pages


# ============================================================================
# 3. DOCX Extractor
# 3. DOCX Extractor
# ============================================================================
def _extract_docx(file_bytes: bytes, filename: str) -> List[PageContent]:
    # DOCX doesn't have fixed pages / DOCX ไม่มีหน้าคงที่
    page_content = PageContent(page_number=1)
    
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        
        pos_counter = 0
        
        # 1. Extract Text Paragraphs / ดึงย่อหน้าข้อความ
        for para in doc.paragraphs:
            if para.text.strip():
                page_content.text_blocks.append(DocumentElement(
                    page_number=1,
                    position_in_page=pos_counter,
                    y_position=pos_counter * 20,  # Fake Y-pos
                    element_type='text',
                    content=para.text.strip()
                ))
                pos_counter += 1
                
        # 2. Extract Tables (as Text) / ดึงตาราง (เป็นข้อความ)
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                table_text.append(row_text)
            
            full_table = "\n".join(table_text)
            if full_table.strip():
                 page_content.text_blocks.append(DocumentElement(
                    page_number=1,
                    position_in_page=pos_counter,
                    y_position=pos_counter * 20,
                    element_type='text',
                    content="[TABLE]\n" + full_table
                ))
                 pos_counter += 1
        
        # 3. Extract Images / ดึงรูปภาพ
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type  # e.g. image/jpeg
                    ext = content_type.split('/')[-1] if '/' in content_type else 'png'
                    
                    page_content.images.append(DocumentElement(
                        page_number=1,
                        position_in_page=pos_counter,
                        y_position=0,
                        element_type='image',
                        image_bytes=image_bytes,
                        image_ext=ext
                    ))
                    pos_counter += 1
                except Exception:
                    pass
        
        print(f"✅ Extracted DOCX content ({len(page_content.text_blocks)} blocks)")
    except Exception as e:
        print(f"❌ Error extracting DOCX: {e}")
        
    return [page_content]


# ============================================================================
# 4. Excel/CSV Extractor
# 4. Excel/CSV Extractor
# ============================================================================
def _extract_excel_csv(file_bytes: bytes, filename: str, ext: str) -> List[PageContent]:
    page_content = PageContent(page_number=1)
    try:
        bio = io.BytesIO(file_bytes)
        dfs = {}
        
        if ext == '.csv':
            dfs['Sheet1'] = pd.read_csv(bio)
        else:
            # Excel / Excel
            excel_file = pd.ExcelFile(bio)
            for sheet_name in excel_file.sheet_names:
                dfs[sheet_name] = excel_file.parse(sheet_name)

        pos = 0
        for sheet_name, df in dfs.items():
            markdown_table = df.to_markdown(index=False)
            header = f"### Sheet: {sheet_name}\n"
            content = header + markdown_table
            
            page_content.text_blocks.append(DocumentElement(
                page_number=1,
                position_in_page=pos,
                y_position=pos * 100,
                element_type='text',
                content=content
            ))
            pos += 1
            
        print(f"✅ Extracted Excel/CSV data")
    except Exception as e:
        print(f"❌ Error extracting Excel/CSV: {e}")
        
    return [page_content]


# ============================================================================
# 5. Image File Extractor
# 5. Image File Extractor
# ============================================================================
def _extract_image_file(file_bytes: bytes, filename: str) -> List[PageContent]:
    # Treat as single page with 1 big image / ถือเป็นหน้าเดียวมีรูปใหญ่ 1 รูป
    page_content = PageContent(page_number=1)
    
    # 1. Set as Full Page Image / ตั้งเป็นรูปเต็มหน้า
    page_content.full_page_image = file_bytes 
    
    # 2. Get extension / ดึงนามสกุล
    try:
        img = PilImage.open(io.BytesIO(file_bytes))
        ext = img.format.lower() if img.format else "png"
    except:
        ext = "png"
        
    # 3. Add as image element / เพิ่มเป็น image element
    page_content.images.append(DocumentElement(
        page_number=1,
        position_in_page=0,
        y_position=0,
        element_type='image',
        image_bytes=file_bytes,
        image_ext=ext
    ))
    
    return [page_content]


# ============================================================================
# 6. Generic Text Extractor (TXT, JSON, XML, HTML, RTF)
# 6. Generic Text Extractor
# ============================================================================
def _extract_text_file(file_bytes: bytes, filename: str) -> List[PageContent]:
    page_content = PageContent(page_number=1)
    
    try:
        # Attempt UTF-8, fallback to latin-1 / ลอง UTF-8 ก่อน ไม่ได้ใช้ latin-1
        try:
            text_content = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            text_content = file_bytes.decode('latin-1')
            
        page_content.text_blocks.append(DocumentElement(
            page_number=1,
            position_in_page=0,
            y_position=0,
            element_type='text',
            content=text_content
        ))
        print(f"✅ Extracted Text content from {filename}")
    except Exception as e:
        print(f"❌ Error extracting Text file: {e}")
        # Add error placeholder / เพิ่ม placeholder ข้อผิดพลาด
        page_content.text_blocks.append(DocumentElement(
            page_number=1,
            position_in_page=0,
            y_position=0,
            element_type='text',
            content=f"[Error decoding file content: {str(e)}]"
        ))
        
    return [page_content]


# ============================================================================
# 7. OpenDocument Extractor (ODT, ODS, ODP)
# 7. OpenDocument Extractor
# ============================================================================
def _extract_odf(file_bytes: bytes, filename: str) -> List[PageContent]:
    page_content = PageContent(page_number=1)
    
    try:
        # Load ODF file / โหลดไฟล์ ODF
        doc = load_odf(io.BytesIO(file_bytes))
        
        # Extract text from paragraphs / ดึงข้อความจากย่อหน้า
        all_text = []
        for element in doc.getElementsByType(text.P):
            all_text.append(teletype.extractText(element))
            
        full_text = "\n".join(all_text)
        
        if full_text.strip():
            page_content.text_blocks.append(DocumentElement(
                page_number=1,
                position_in_page=0,
                y_position=0,
                element_type='text',
                content=full_text
            ))
            
        print(f"✅ Extracted ODF content from {filename}")
    except Exception as e:
        print(f"❌ Error extracting ODF: {e}")
        page_content.text_blocks.append(DocumentElement(
            page_number=1,
            position_in_page=0,
            y_position=0,
            element_type='text',
            content=f"[Error extracting ODF content: {str(e)}]"
        ))

    return [page_content]


# ============================================================================
# 8. ZIP Extractor
# 8. ZIP Extractor
# ============================================================================
def _extract_zip(file_bytes: bytes, filename: str) -> List[PageContent]:
    page_content = PageContent(page_number=1)
    
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes), 'r') as zip_ref:
            file_list = zip_ref.namelist()
            
            # 1. List files / แสดงรายชื่อไฟล์
            content_summary = "### ZIP Contents:\n" + "\n".join([f"- {f}" for f in file_list])
            page_content.text_blocks.append(DocumentElement(
                 page_number=1,
                 position_in_page=0,
                 y_position=0,
                 element_type='text',
                 content=content_summary
            ))
            
            # 2. Extract text from internal text files / ดึงข้อความจากไฟล์ข้อความภายใน
            for internal_file in file_list:
                _, ext = os.path.splitext(internal_file)
                if ext.lower() in ['.txt', '.json', '.xml', '.md', '.log']:
                    try:
                        with zip_ref.open(internal_file) as f:
                            text_data = f.read().decode('utf-8')[:5000]  # Limit to 5KB
                            page_content.text_blocks.append(DocumentElement(
                                page_number=1,
                                position_in_page=1,
                                y_position=100,
                                element_type='text',
                                content=f"\n--- File: {internal_file} ---\n{text_data}"
                            ))
                    except:
                        pass
                        
        print(f"✅ Extracted ZIP content from {filename}")
    except Exception as e:
        print(f"❌ Error extracting ZIP: {e}")
        page_content.text_blocks.append(DocumentElement(
            page_number=1,
            position_in_page=0,
            y_position=0,
            element_type='text',
            content=f"[Error parsing ZIP: {str(e)}]"
        ))

    return [page_content]


# ============================================================================
# 9. TTF Extractor (Metadata)
# 9. TTF Extractor (Metadata)
# ============================================================================
def _extract_ttf(file_bytes: bytes, filename: str) -> List[PageContent]:
    page_content = PageContent(page_number=1)
    
    try:
        font = TTFont(io.BytesIO(file_bytes))
        
        # Name record IDs: 0=Copyright, 1=FontFamily, 4=FullFontName
        name_records = ""
        for record in font['name'].names:
            if record.nameID in [0, 1, 3, 4]:
                try:
                    s = record.string.decode('utf-16-be' if b'\x00' in record.string else 'latin-1')
                    label = {0: "Copyright", 1: "Family", 3: "Unique ID", 4: "Full Name"}.get(record.nameID, "Unknown")
                    name_records += f"- {label}: {s}\n"
                except:
                    pass
                    
        content = f"### Font Metadata ({filename})\n{name_records}"
        
        page_content.text_blocks.append(DocumentElement(
            page_number=1,
            position_in_page=0,
            y_position=0,
            element_type='text',
            content=content
        ))
        print(f"✅ Extracted TTF metadata from {filename}")
    except Exception as e:
        print(f"❌ Error extracting TTF: {e}")
        page_content.text_blocks.append(DocumentElement(
            page_number=1,
            position_in_page=0,
            y_position=0,
            element_type='text',
            content=f"[Error parsing TTF: {str(e)}]"
        ))

    return [page_content]
